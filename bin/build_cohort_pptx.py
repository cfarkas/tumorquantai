#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
bin/build_cohort_pptx.py

Cancer-agnostic compact single-deck PowerPoint builder for LazySlide / HistoPLUS outputs.

Why V3 exists
-------------
The original script could create one PPTX for one disease-specific cohort, but the deck text
and example paths were lymphoma-oriented. V3 makes the report cancer-agnostic and adds automatic
cohort-level sample x cell-type CSV outputs.

Expected folder structure
-------------------------
ROOT/
  case_001/
    cell_types/class_counts.csv
    overlays/celltypes_overview_and_zoom.pdf
    overlays/celltypes_overview_and_zoom.png
    post_visualization/spatial_rois_report.pdf
    post_visualization/cell_embedding_umap_report.pdf          optional
    post_visualization/spatial_rois/                           optional; searched recursively
    summary/summary.json                                       optional
    summary/run_metadata.json                                  optional

Main new flags
--------------
  --compact-first-pages
      Compact mode. Keeps only the first page of spatial_rois_report.pdf and, when --include-umap
      is used, the first page of cell_embedding_umap_report.pdf. It also enables spotlight search
      unless --no-spotlights is passed.

  --include-spotlights
      Adds class-specific whole-slide H&E spotlight material. The script searches recursively for
      image/PDF files whose paths contain terms like spotlight, whole_slide, class_specific, he,
      h_e, or wsi. It also searches spatial_rois_report.pdf page text for spotlight-like terms.

  --spotlights-per-slide 4
      Places spotlight images in a compact grid, up to 4 per slide, to avoid generating too many
      slides.

Accompanying cell-stat CSV outputs
----------------------------------
By default, every run writes two cohort-level matrices next to the PPTX:

  <PPTX>.cell_counts_by_sample.csv
      Rows are detected samples and columns are detected cell types/classes. Values are raw
      detected cell counts from each sample's cell_types/class_counts.csv.

  <PPTX>.cell_fractions_by_sample.csv
      Same shape, but each row is normalized by all detected cells in that sample. Values are
      fractions from 0 to 1 and each non-empty sample row sums to 1.

Use --cell-stats-csv and --cell-stats-normalized-csv to choose explicit output paths, or
--no-cell-stats-csv to disable these companion tables.

Install dependencies
--------------------
conda activate lazyslide311
pip install python-pptx pymupdf pillow pandas matplotlib

Recommended compact command
---------------------------
python bin/build_cohort_pptx.py \
  --root /path/to/histoplus_results \
  --out /path/to/histoplus_results/pptx_reports/cohort_report.pptx \
  --dpi 180 \
  --include-umap \
  --compact-first-pages \
  --clean-cache \
  --force-render

Verify new flags
----------------
python bin/build_cohort_pptx.py --help | grep -E "compact|spotlight|cell-stats|section"
"""

from __future__ import annotations

import argparse
import fnmatch
import json
import math
import re
import shutil
import sys
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

SCRIPT_VERSION = "V3.0.0-cancer-agnostic-cell-stats"


def _die_missing(package: str, install_hint: str) -> None:
    print(f"ERROR: missing Python package: {package}", file=sys.stderr)
    print(f"Install with: {install_hint}", file=sys.stderr)
    sys.exit(2)


try:
    import fitz  # PyMuPDF
except Exception:
    _die_missing("pymupdf", "pip install pymupdf")

try:
    import pandas as pd
except Exception:
    _die_missing("pandas", "pip install pandas")

try:
    from PIL import Image
except Exception:
    _die_missing("pillow", "pip install pillow")

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
except Exception:
    _die_missing("matplotlib", "pip install matplotlib")

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception:
    _die_missing("python-pptx", "pip install python-pptx")


Image.MAX_IMAGE_PIXELS = None

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_X = 0.42
TITLE_Y = 0.22
TITLE_H = 0.45
FOOTER_Y = 7.18

COLOR_DARK = RGBColor(30, 33, 38)
COLOR_MUTED = RGBColor(88, 94, 104)
COLOR_LIGHT = RGBColor(246, 247, 249)
COLOR_LINE = RGBColor(214, 218, 224)
COLOR_BLUE = RGBColor(35, 86, 160)
COLOR_GREEN = RGBColor(32, 139, 90)
COLOR_ORANGE = RGBColor(194, 112, 24)
COLOR_RED = RGBColor(184, 61, 61)
COLOR_WHITE = RGBColor(255, 255, 255)

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp", ".bmp"}
PDF_SUFFIXES = {".pdf"}
SPOTLIGHT_SUFFIXES = IMAGE_SUFFIXES | PDF_SUFFIXES

DEFAULT_SPOTLIGHT_TERMS = [
    "spotlight",
    "whole_slide",
    "whole-slide",
    "whole slide",
    "whole wsi",
    "wsi",
    "class_specific",
    "class-specific",
    "class specific",
    "he_spotlight",
    "h&e",
    "h_e",
    "he_only",
    "he-only",
    "hematoxylin",
]

EXCLUDE_SPOTLIGHT_STEMS = {
    "spatial_rois_report",
    "cell_embedding_umap_report",
    "celltypes_overview_and_zoom",
    "celltype_overview",
    "celltypes_overview",
    "cell_type_overview",
    "overview_with_zoom_box",
    "zoom_overlay_celltypes",
}


@dataclass
class CaseAssets:
    name: str
    path: Path
    overview_pdf: Optional[Path]
    overview_png: Optional[Path]
    roi_pdf: Optional[Path]
    umap_pdf: Optional[Path]
    counts_csv: Optional[Path]
    summary_json: Optional[Path]
    run_metadata_json: Optional[Path]


@dataclass
class BuildStats:
    out: Path
    cases: int = 0
    slides: int = 0
    overview_slides: int = 0
    count_slides: int = 0
    roi_slides: int = 0
    umap_slides: int = 0
    spotlight_slides: int = 0
    spotlight_items: int = 0
    spotlight_files: int = 0
    roi_spotlight_pages: int = 0
    missing_overview: List[str] = field(default_factory=list)
    missing_roi: List[str] = field(default_factory=list)
    missing_counts: List[str] = field(default_factory=list)
    missing_spotlights: List[str] = field(default_factory=list)
    case_records: List[Dict[str, Any]] = field(default_factory=list)

    def as_dict(self) -> Dict[str, Any]:
        return {
            "out": str(self.out),
            "cases": self.cases,
            "slides": self.slides,
            "overview_slides": self.overview_slides,
            "count_slides": self.count_slides,
            "roi_slides": self.roi_slides,
            "umap_slides": self.umap_slides,
            "spotlight_slides": self.spotlight_slides,
            "spotlight_items": self.spotlight_items,
            "spotlight_files": self.spotlight_files,
            "roi_spotlight_pages": self.roi_spotlight_pages,
            "missing_overview": self.missing_overview,
            "missing_roi": self.missing_roi,
            "missing_counts": self.missing_counts,
            "missing_spotlights": self.missing_spotlights,
            "case_records": self.case_records,
        }


def natural_key(text: str) -> List[Any]:
    parts = re.split(r"(\d+)", str(text))
    out: List[Any] = []
    for part in parts:
        if part.isdigit():
            out.append(int(part))
        else:
            out.append(part.lower())
    return out


def slugify(text: str, max_len: int = 120) -> str:
    text = re.sub(r"[^A-Za-z0-9._-]+", "_", str(text)).strip("._-")
    text = re.sub(r"_+", "_", text)
    return (text or "case")[:max_len]


def safe_exists(path: Optional[Path]) -> bool:
    return bool(path and path.exists() and path.is_file() and path.stat().st_size > 0)


def first_existing(paths: Sequence[Path]) -> Optional[Path]:
    for p in paths:
        if safe_exists(p):
            return p
    return None


def rel_path(path: Optional[Path], root: Path) -> str:
    if path is None:
        return "missing"
    try:
        return str(path.relative_to(root))
    except Exception:
        return str(path)


def collect_case_assets(case_dir: Path) -> Optional[CaseAssets]:
    overview_pdf = first_existing([
        case_dir / "overlays" / "celltypes_overview_and_zoom.pdf",
        case_dir / "overlays" / "celltype_overview.pdf",
        case_dir / "overlays" / "celltypes_overview.pdf",
        case_dir / "overlays" / "cell_type_overview.pdf",
    ])
    overview_png = first_existing([
        case_dir / "overlays" / "celltypes_overview_and_zoom.png",
        case_dir / "overlays" / "overview_with_zoom_box.png",
        case_dir / "overlays" / "zoom_overlay_celltypes.png",
    ])
    roi_pdf = first_existing([
        case_dir / "post_visualization" / "spatial_rois_report.pdf",
    ])
    umap_pdf = first_existing([
        case_dir / "post_visualization" / "cell_embedding_umap_report.pdf",
    ])
    counts_csv = first_existing([
        case_dir / "cell_types" / "class_counts.csv",
    ])
    summary_json = first_existing([
        case_dir / "summary" / "summary.json",
        case_dir / "post_visualization" / "post_visualization_summary.json",
    ])
    run_metadata_json = first_existing([
        case_dir / "summary" / "run_metadata.json",
    ])

    if not any([safe_exists(overview_pdf), safe_exists(overview_png), safe_exists(roi_pdf), safe_exists(umap_pdf), safe_exists(counts_csv)]):
        return None

    return CaseAssets(
        name=case_dir.name,
        path=case_dir,
        overview_pdf=overview_pdf,
        overview_png=overview_png,
        roi_pdf=roi_pdf,
        umap_pdf=umap_pdf,
        counts_csv=counts_csv,
        summary_json=summary_json,
        run_metadata_json=run_metadata_json,
    )


def discover_cases(root: Path, recursive: bool, include: Sequence[str], exclude: Sequence[str], first_n: int = 0) -> List[CaseAssets]:
    if recursive:
        candidate_dirs = [p for p in root.rglob("*") if p.is_dir()]
    else:
        candidate_dirs = [p for p in root.iterdir() if p.is_dir()]

    cases: List[CaseAssets] = []
    seen: set[Path] = set()
    ignored_names = {
        "pptx_reports", "per_case", "working", "summary", "cell_types", "overlays",
        "post_visualization", "qc_patches", "spatial_rois", "_render_cache_all_samples",
    }
    for d in sorted(candidate_dirs, key=lambda p: natural_key(p.name)):
        if d in seen:
            continue
        if d.name.startswith(".") or d.name.startswith("_"):
            continue
        if d.name in ignored_names:
            continue
        if include and not any(fnmatch.fnmatch(d.name, pat) for pat in include):
            continue
        if exclude and any(fnmatch.fnmatch(d.name, pat) for pat in exclude):
            continue
        case = collect_case_assets(d)
        if case is None:
            continue
        if recursive and any(parent in seen for parent in d.parents):
            continue
        cases.append(case)
        seen.add(d)
        if first_n > 0 and len(cases) >= first_n:
            break
    return cases


def read_json(path: Optional[Path]) -> Dict[str, Any]:
    if not safe_exists(path):
        return {}
    try:
        with open(path, "r", encoding="utf-8") as f:
            obj = json.load(f)
        return obj if isinstance(obj, dict) else {}
    except Exception:
        return {}


def infer_counts_table(csv_path: Path) -> Optional[pd.DataFrame]:
    try:
        df = pd.read_csv(csv_path)
    except Exception:
        return None
    if df.empty:
        return None

    normalized = {c: re.sub(r"[^a-z0-9]+", "_", str(c).lower()).strip("_") for c in df.columns}
    name_candidates = [
        "class", "classes", "cell_type", "celltype", "celltypes", "type", "label", "labels",
        "prediction", "predicted_class", "class_name", "name"
    ]
    count_candidates = [
        "count", "counts", "n", "num", "number", "freq", "frequency", "cells", "n_cells",
        "cell_count", "total", "total_cells"
    ]

    name_col = None
    for wanted in name_candidates:
        for col, norm in normalized.items():
            if norm == wanted:
                name_col = col
                break
        if name_col is not None:
            break
    if name_col is None:
        object_cols = [c for c in df.columns if df[c].dtype == object]
        name_col = object_cols[0] if object_cols else df.columns[0]

    count_col = None
    for wanted in count_candidates:
        for col, norm in normalized.items():
            if col == name_col:
                continue
            if norm == wanted:
                count_col = col
                break
        if count_col is not None:
            break
    if count_col is None:
        for col in df.columns:
            if col == name_col:
                continue
            numeric = pd.to_numeric(df[col], errors="coerce")
            if numeric.notna().sum() > 0:
                df[col] = numeric
                count_col = col
                break
    if count_col is None:
        return None

    out = df[[name_col, count_col]].copy()
    out.columns = ["class", "count"]
    out["class"] = out["class"].astype(str)
    out["count"] = pd.to_numeric(out["count"], errors="coerce")
    out = out.dropna(subset=["class", "count"])
    out = out[out["count"] > 0]
    if out.empty:
        return None
    out = out.groupby("class", as_index=False)["count"].sum().sort_values("count", ascending=False)
    total = out["count"].sum()
    out["percent"] = 100.0 * out["count"] / total if total else 0.0
    return out


def read_case_cell_total(case: CaseAssets) -> Optional[int]:
    if not safe_exists(case.counts_csv):
        return None
    tbl = infer_counts_table(case.counts_csv)  # type: ignore[arg-type]
    if tbl is None or tbl.empty:
        return None
    return int(tbl["count"].sum())


def aggregate_counts(cases: Sequence[CaseAssets]) -> Optional[pd.DataFrame]:
    frames = []
    for case in cases:
        if not safe_exists(case.counts_csv):
            continue
        tbl = infer_counts_table(case.counts_csv)  # type: ignore[arg-type]
        if tbl is None or tbl.empty:
            continue
        tbl = tbl[["class", "count"]].copy()
        tbl["sample"] = case.name
        frames.append(tbl)
    if not frames:
        return None
    all_counts = pd.concat(frames, ignore_index=True)
    agg = all_counts.groupby("class", as_index=False)["count"].sum().sort_values("count", ascending=False)
    total = agg["count"].sum()
    agg["percent"] = 100.0 * agg["count"] / total if total else 0.0
    return agg




def build_wide_cell_stats_tables(cases: Sequence[CaseAssets]) -> Tuple[pd.DataFrame, pd.DataFrame, Dict[str, Any]]:
    """Build sample x cell-type matrices from each case's cell_types/class_counts.csv.

    Returns
    -------
    counts_wide:
        Rows are samples and columns are cell classes/types. Values are raw detected cell counts.
    fractions_wide:
        Same shape, row-normalized by total detected cells in each sample. Values are fractions
        from 0 to 1. Rows with zero detected cells are filled with 0.
    info:
        Small summary dictionary for logging/manifest output.
    """
    sample_names = [case.name for case in cases]
    count_maps: Dict[str, Dict[str, int]] = {name: {} for name in sample_names}
    class_totals: Dict[str, int] = {}
    missing_counts: List[str] = []
    unreadable_counts: List[str] = []

    for case in cases:
        if not safe_exists(case.counts_csv):
            missing_counts.append(case.name)
            continue

        tbl = infer_counts_table(case.counts_csv)  # type: ignore[arg-type]
        if tbl is None or tbl.empty:
            unreadable_counts.append(case.name)
            continue

        for _, row in tbl.iterrows():
            cls = str(row["class"]).strip()
            if not cls:
                continue
            count = int(round(float(row["count"])))
            if count < 0:
                continue
            count_maps[case.name][cls] = count_maps[case.name].get(cls, 0) + count
            class_totals[cls] = class_totals.get(cls, 0) + count

    ordered_classes = sorted(class_totals, key=lambda c: (-class_totals[c], natural_key(c)))
    counts_wide = pd.DataFrame(0, index=sample_names, columns=ordered_classes, dtype="int64")
    counts_wide.index.name = "sample"

    for sample, cmap in count_maps.items():
        for cls, count in cmap.items():
            if cls in counts_wide.columns:
                counts_wide.loc[sample, cls] = int(count)

    row_totals = counts_wide.sum(axis=1)
    denom = row_totals.replace(0, pd.NA)
    fractions_wide = counts_wide.astype(float).div(denom, axis=0).fillna(0.0)
    fractions_wide.index.name = "sample"

    info = {
        "samples": len(sample_names),
        "samples_with_counts": int((row_totals > 0).sum()),
        "detected_cell_types": len(ordered_classes),
        "total_detected_cells": int(row_totals.sum()),
        "missing_counts": missing_counts,
        "unreadable_counts": unreadable_counts,
        "columns_ordered_by": "descending total detected cells across all samples",
        "normalized_values": "fraction of all detected cells in each sample; non-empty rows sum to 1",
    }
    return counts_wide, fractions_wide, info


def write_cell_stats_csvs(cases: Sequence[CaseAssets], raw_csv: Path, normalized_csv: Path) -> Dict[str, Any]:
    """Write companion raw and normalized sample x cell-type CSV matrices."""
    raw_csv.parent.mkdir(parents=True, exist_ok=True)
    normalized_csv.parent.mkdir(parents=True, exist_ok=True)
    counts_wide, fractions_wide, info = build_wide_cell_stats_tables(cases)
    counts_wide.to_csv(raw_csv)
    fractions_wide.to_csv(normalized_csv, float_format="%.10g")
    info = dict(info)
    info["raw_counts_csv"] = str(raw_csv)
    info["normalized_fractions_csv"] = str(normalized_csv)
    return info

def render_barh(df: pd.DataFrame, out_png: Path, title: str, top_n: int = 15) -> Path:
    top = df.head(top_n).copy()
    top["label"] = top["class"].astype(str).map(lambda s: s if len(s) <= 34 else s[:31] + "...")
    h = max(3.4, min(7.0, 0.34 * len(top) + 1.15))
    fig = plt.figure(figsize=(8.6, h))
    ax = fig.add_subplot(111)
    ax.barh(top["label"].iloc[::-1], top["count"].iloc[::-1])
    ax.set_xlabel("Cell count")
    ax.set_title(title)
    ax.grid(axis="x", alpha=0.22)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    fig.tight_layout()
    out_png.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(str(out_png), dpi=180, bbox_inches="tight")
    plt.close(fig)
    return out_png


def pdf_page_count(pdf_path: Optional[Path]) -> int:
    if not safe_exists(pdf_path):
        return 0
    try:
        doc = fitz.open(str(pdf_path))  # type: ignore[arg-type]
        n = len(doc)
        doc.close()
        return n
    except Exception:
        return 0


def render_pdf_pages(pdf_path: Path, out_dir: Path, dpi: int, max_pages: Optional[int], force: bool) -> List[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf_path))
    n_pages = len(doc)
    if n_pages == 0:
        doc.close()
        return []
    if max_pages is not None:
        n_pages = min(n_pages, max_pages)
    indices = list(range(n_pages))
    doc.close()
    return render_pdf_page_indices(pdf_path, out_dir, dpi=dpi, page_indices=indices, force=force)


def render_pdf_page_indices(pdf_path: Path, out_dir: Path, dpi: int, page_indices: Sequence[int], force: bool) -> List[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    if not page_indices:
        return []
    doc = fitz.open(str(pdf_path))
    mat = fitz.Matrix(dpi / 72.0, dpi / 72.0)
    rendered: List[Path] = []
    stem = slugify(pdf_path.stem)
    for idx in page_indices:
        if idx < 0 or idx >= len(doc):
            continue
        out = out_dir / f"{stem}_p{idx + 1:03d}_dpi{dpi}.png"
        if safe_exists(out) and not force:
            rendered.append(out)
            continue
        page = doc.load_page(idx)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        pix.save(str(out))
        rendered.append(out)
    doc.close()
    return rendered


def normalize_terms(terms: Sequence[str]) -> List[str]:
    out: List[str] = []
    for term in terms:
        term = str(term).strip().lower()
        if not term:
            continue
        out.append(term)
        out.append(term.replace("-", "_"))
        out.append(term.replace(" ", "_"))
    # keep order while deduplicating
    seen: set[str] = set()
    dedup: List[str] = []
    for term in out:
        if term not in seen:
            dedup.append(term)
            seen.add(term)
    return dedup


def _candidate_text(path: Path, case_dir: Path) -> str:
    try:
        rel = str(path.relative_to(case_dir)).lower()
    except Exception:
        rel = str(path).lower()
    rel = rel.replace("\\", "/")
    compact = rel.replace("-", "_").replace(" ", "_")
    return rel + "\n" + compact


def is_spotlight_candidate(path: Path, case_dir: Path, terms: Sequence[str], include_all_spatial_roi_images: bool) -> bool:
    if not path.is_file() or path.suffix.lower() not in SPOTLIGHT_SUFFIXES:
        return False
    if path.stem.lower() in EXCLUDE_SPOTLIGHT_STEMS:
        return False

    text = _candidate_text(path, case_dir)

    # Strict candidate: file/path contains spotlight-like terms.
    normalized = normalize_terms(terms)
    if any(term in text for term in normalized):
        # Avoid adding generic non-H&E/overview files just because they contain WSI in parent paths.
        return True

    # Optional broad mode: include all images/PDFs inside post_visualization/spatial_rois.
    if include_all_spatial_roi_images:
        try:
            rel_parts = path.relative_to(case_dir).parts
            rel_lower = [p.lower() for p in rel_parts]
            if "post_visualization" in rel_lower and "spatial_rois" in rel_lower:
                return True
        except Exception:
            pass

    return False


def discover_spotlight_files(
    case: CaseAssets,
    terms: Sequence[str],
    extra_globs: Sequence[str],
    include_all_spatial_roi_images: bool,
    max_files: int,
) -> List[Path]:
    roots = [
        case.path / "post_visualization" / "spatial_rois",
        case.path / "post_visualization",
        case.path / "overlays",
    ]

    found: Dict[Path, Path] = {}
    for root in roots:
        if not root.exists() or not root.is_dir():
            continue
        for path in root.rglob("*"):
            if is_spotlight_candidate(path, case.path, terms, include_all_spatial_roi_images):
                found[path.resolve()] = path

    for pattern in extra_globs:
        for path in case.path.glob(pattern):
            if path.is_file() and path.suffix.lower() in SPOTLIGHT_SUFFIXES:
                if path.stem.lower() not in EXCLUDE_SPOTLIGHT_STEMS:
                    found[path.resolve()] = path

    files = list(found.values())
    files.sort(key=lambda p: natural_key(str(p.relative_to(case.path)) if p.is_relative_to(case.path) else str(p)))
    if max_files > 0:
        files = files[:max_files]
    return files


def find_pdf_pages_by_terms(pdf_path: Optional[Path], terms: Sequence[str], exclude_indices: Iterable[int], max_pages: int) -> List[int]:
    if not safe_exists(pdf_path):
        return []
    normalized = normalize_terms(terms)
    exclude = set(exclude_indices)
    matched: List[int] = []
    try:
        doc = fitz.open(str(pdf_path))  # type: ignore[arg-type]
        for idx in range(len(doc)):
            if idx in exclude:
                continue
            try:
                text = doc.load_page(idx).get_text("text").lower()
            except Exception:
                text = ""
            text_compact = text.replace("-", "_").replace(" ", "_")
            if any(term in text or term in text_compact for term in normalized):
                matched.append(idx)
                if max_pages > 0 and len(matched) >= max_pages:
                    break
        doc.close()
    except Exception:
        return []
    return matched


def image_size(path: Path) -> Tuple[int, int]:
    with Image.open(path) as im:
        return im.size


def ensure_pptx_image(image_path: Path, out_dir: Path, force: bool = False) -> Path:
    """Return a PNG/JPEG-compatible image path for python-pptx."""
    suffix = image_path.suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg"}:
        return image_path
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / f"{slugify(image_path.stem)}.png"
    if safe_exists(out) and not force:
        return out
    with Image.open(image_path) as im:
        if im.mode not in {"RGB", "RGBA"}:
            im = im.convert("RGB")
        if im.mode == "RGBA":
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=im.split()[-1])
            im = bg
        im.save(out)
    return out


def set_bg(slide, color: RGBColor = COLOR_WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: float = 12,
    bold: bool = False,
    color: RGBColor = COLOR_DARK,
    align: Optional[int] = None,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = str(text).split("\n") if text is not None else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str = "") -> None:
    add_text(slide, title, MARGIN_X, TITLE_Y, 10.8, TITLE_H, font_size=20, bold=True, color=COLOR_DARK)
    if subtitle:
        add_text(slide, subtitle, 9.2, TITLE_Y + 0.05, 3.7, 0.28, font_size=8.5, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_X), Inches(0.72), Inches(12.5), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LINE
    line.line.color.rgb = COLOR_LINE


def add_footer(slide, sample: str, source: str = "") -> None:
    txt = sample if not source else f"{sample}  |  {source}"
    add_text(slide, txt, MARGIN_X, FOOTER_Y, 12.45, 0.2, font_size=7.2, color=COLOR_MUTED)


def add_badge(slide, text: str, x: float, y: float, w: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    add_text(slide, text, x + 0.05, y + 0.045, w - 0.1, 0.15, font_size=7.0, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


def add_image_contain(slide, image_path: Path, x: float, y: float, w: float, h: float, border: bool = True) -> None:
    iw, ih = image_size(image_path)
    if iw <= 0 or ih <= 0:
        raise ValueError(f"Invalid image dimensions: {image_path}")
    img_ratio = iw / ih
    box_ratio = w / h
    if img_ratio >= box_ratio:
        new_w = w
        new_h = w / img_ratio
        new_x = x
        new_y = y + (h - new_h) / 2
    else:
        new_h = h
        new_w = h * img_ratio
        new_x = x + (w - new_w) / 2
        new_y = y
    if border:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_WHITE
        bg.line.color.rgb = COLOR_LINE
        bg.line.width = Pt(0.7)
    slide.shapes.add_picture(str(image_path), Inches(new_x), Inches(new_y), width=Inches(new_w), height=Inches(new_h))


def short_label_from_path(path: Path, case_dir: Path, max_len: int = 46) -> str:
    try:
        rel = path.relative_to(case_dir)
        label = str(rel.with_suffix(""))
    except Exception:
        label = path.stem
    label = label.replace("post_visualization/spatial_rois/", "")
    label = label.replace("post_visualization/", "")
    label = label.replace("overlays/", "")
    label = label.replace("_", " ").replace("-", "-")
    label = re.sub(r"\s+", " ", label).strip()
    if len(label) > max_len:
        label = "…" + label[-(max_len - 1):]
    return label


def chunked(items: Sequence[Any], n: int) -> Iterable[Sequence[Any]]:
    for i in range(0, len(items), n):
        yield items[i:i+n]


def add_title_slide(prs: Presentation, cases: Sequence[CaseAssets], root: Path, include_umap: bool, compact: bool, include_spotlights: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_WHITE)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_DARK
    band.line.color.rgb = COLOR_DARK

    add_text(slide, "Cancer-agnostic spatial cell-type reports", 0.75, 2.05, 11.85, 0.65, font_size=30, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    mode = "compact first-page deck" if compact else "single consolidated PPTX"
    add_text(slide, f"{mode} | {len(cases)} samples | {SCRIPT_VERSION}", 1.2, 2.92, 10.9, 0.40, font_size=15.5, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    subtitle = "LazySlide / HistoPLUS output: class counts, overview/zoom, spatial ROIs"
    if include_umap:
        subtitle += ", UMAP"
    if include_spotlights:
        subtitle += ", and class-specific whole-slide H&E spotlights"
    add_text(slide, subtitle, 1.2, 3.52, 10.9, 0.35, font_size=11.5, color=COLOR_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, f"Root: {root}", 0.85, 6.82, 11.6, 0.25, font_size=8, color=COLOR_MUTED, align=PP_ALIGN.CENTER)


def add_dataset_summary_slide(prs: Presentation, cases: Sequence[CaseAssets], root: Path, args: argparse.Namespace) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Dataset overview", f"{len(cases)} samples detected")

    cols = ["Sample", "cells", "overview", "ROI pages", "UMAP pages", "spotlight files", "counts"]
    rows = []
    for case in cases:
        cells = read_case_cell_total(case)
        spot_files = discover_spotlight_files(
            case,
            terms=args.spotlight_terms,
            extra_globs=args.spotlight_globs,
            include_all_spatial_roi_images=args.include_all_spatial_roi_images,
            max_files=args.max_spotlight_pages,
        ) if args.include_spotlights else []
        rows.append([
            case.name,
            f"{cells:,}" if cells is not None else "—",
            "yes" if (safe_exists(case.overview_pdf) or safe_exists(case.overview_png)) else "no",
            str(pdf_page_count(case.roi_pdf)) if safe_exists(case.roi_pdf) else "0",
            str(pdf_page_count(case.umap_pdf)) if safe_exists(case.umap_pdf) else "0",
            str(len(spot_files)) if args.include_spotlights else "off",
            "yes" if safe_exists(case.counts_csv) else "no",
        ])

    table = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(0.42), Inches(0.95), Inches(12.5), Inches(5.95)).table
    widths = [2.78, 1.10, 1.05, 1.05, 1.05, 1.18, 0.95]
    total_width = sum(widths)
    scale = 12.5 / total_width
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w * scale)

    for j, col in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(8.0)
                run.font.bold = True
                run.font.color.rgb = COLOR_WHITE
                run.font.name = "Aptos"

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(6.9)
                    run.font.color.rgb = COLOR_DARK
                    run.font.name = "Aptos"

    add_footer(slide, "all samples", str(root))


def add_aggregate_counts_slide(prs: Presentation, cases: Sequence[CaseAssets], cache_dir: Path, top_classes: int) -> int:
    agg = aggregate_counts(cases)
    if agg is None or agg.empty:
        return 0
    total_cells = int(agg["count"].sum())
    chart = render_barh(agg, cache_dir / "aggregate_counts" / "aggregate_top_classes.png", f"Top {min(top_classes, len(agg))} cell classes across all samples | total cells: {total_cells:,}", top_n=top_classes)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Aggregate cell-type composition", "all samples")
    add_image_contain(slide, chart, 0.55, 1.0, 8.3, 5.9, border=True)

    lines = [
        f"Samples included: {len(cases)}",
        f"Total cells: {total_cells:,}",
        f"Detected classes: {len(agg)}",
        "",
        "Top classes:",
    ]
    for _, row in agg.head(8).iterrows():
        lines.append(f"• {row['class']}: {int(row['count']):,} ({row['percent']:.1f}%)")
    add_text(slide, "Combined quantitative summary", 9.1, 1.08, 3.6, 0.35, font_size=13.5, bold=True, color=COLOR_DARK)
    add_text(slide, "\n".join(lines), 9.1, 1.5, 3.7, 4.9, font_size=10, color=COLOR_DARK)
    add_footer(slide, "all samples", "cell_types/class_counts.csv merged across cases")
    return 1


def add_section_slide(prs: Presentation, case: CaseAssets, idx: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_LIGHT)
    total_cells = read_case_cell_total(case)
    add_text(slide, f"Sample {idx}/{total}", 0.9, 2.25, 11.5, 0.28, font_size=12.5, bold=True, color=COLOR_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, case.name, 0.9, 2.68, 11.5, 0.65, font_size=30, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    sub = []
    if total_cells is not None:
        sub.append(f"{total_cells:,} cells")
    if safe_exists(case.roi_pdf):
        sub.append(f"{pdf_page_count(case.roi_pdf)} spatial ROI pages available")
    if safe_exists(case.umap_pdf):
        sub.append(f"{pdf_page_count(case.umap_pdf)} UMAP pages available")
    add_text(slide, " | ".join(sub) if sub else str(case.path), 0.9, 3.45, 11.5, 0.32, font_size=10, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, str(case.path), 0.9, 6.85, 11.5, 0.23, font_size=7.5, color=COLOR_MUTED, align=PP_ALIGN.CENTER)


def add_case_counts_slide(prs: Presentation, case: CaseAssets, root: Path, cache_dir: Path, top_classes: int) -> int:
    if not safe_exists(case.counts_csv):
        return 0
    tbl = infer_counts_table(case.counts_csv)  # type: ignore[arg-type]
    if tbl is None or tbl.empty:
        return 0
    total_cells = int(tbl["count"].sum())
    chart = render_barh(tbl, cache_dir / "class_counts" / "top_classes.png", f"Top {min(top_classes, len(tbl))} cell classes | total cells: {total_cells:,}", top_n=top_classes)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, f"{case.name}: cell-type composition", "class_counts.csv")
    add_image_contain(slide, chart, 0.55, 1.0, 8.1, 5.9, border=True)

    lines = [
        f"Total cells: {total_cells:,}",
        f"Detected classes: {len(tbl)}",
        "",
        "Top classes:",
    ]
    for _, row in tbl.head(8).iterrows():
        lines.append(f"• {row['class']}: {int(row['count']):,} ({row['percent']:.1f}%)")
    add_text(slide, "Quantitative summary", 8.95, 1.1, 3.7, 0.35, font_size=13.5, bold=True, color=COLOR_DARK)
    add_text(slide, "\n".join(lines), 8.95, 1.55, 3.75, 4.9, font_size=10, color=COLOR_DARK)
    add_footer(slide, case.name, rel_path(case.counts_csv, root))
    return 1


def add_overview_slide(prs: Presentation, case: CaseAssets, root: Path, cache_dir: Path, dpi: int, overview_source: str, force_render: bool) -> int:
    images: List[Tuple[Path, str]] = []
    if overview_source in {"auto", "pdf"} and safe_exists(case.overview_pdf):
        rendered = render_pdf_pages(case.overview_pdf, cache_dir / "overview", dpi, max_pages=1, force=force_render)  # type: ignore[arg-type]
        images.extend([(p, rel_path(case.overview_pdf, root)) for p in rendered])
    if not images and overview_source in {"auto", "png"} and safe_exists(case.overview_png):
        images.append((case.overview_png, rel_path(case.overview_png, root)))  # type: ignore[arg-type]
    if not images:
        return 0

    n = 0
    for img, source in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_title(slide, f"{case.name}: cell-type overview and zoom", "overview")
        add_image_contain(slide, img, 0.55, 0.92, 12.25, 5.95, border=True)
        add_footer(slide, case.name, source)
        n += 1
    return n


def add_pdf_report_slides(
    prs: Presentation,
    case: CaseAssets,
    pdf_path: Optional[Path],
    root: Path,
    cache_dir: Path,
    dpi: int,
    max_pages: Optional[int],
    force_render: bool,
    title_prefix: str,
    section_label: str,
) -> int:
    if not safe_exists(pdf_path):
        return 0
    rendered = render_pdf_pages(pdf_path, cache_dir / slugify(pdf_path.stem), dpi=dpi, max_pages=max_pages, force=force_render)  # type: ignore[arg-type]
    for i, img in enumerate(rendered, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_title(slide, f"{case.name}: {title_prefix}", f"{section_label} page {i}/{len(rendered)}")
        add_image_contain(slide, img, 0.55, 0.92, 12.25, 5.95, border=True)
        add_footer(slide, case.name, rel_path(pdf_path, root))
    return len(rendered)


def add_pdf_selected_page_slides(
    prs: Presentation,
    case: CaseAssets,
    pdf_path: Optional[Path],
    page_indices: Sequence[int],
    root: Path,
    cache_dir: Path,
    dpi: int,
    force_render: bool,
    title_prefix: str,
    section_label: str,
) -> int:
    if not safe_exists(pdf_path) or not page_indices:
        return 0
    rendered = render_pdf_page_indices(pdf_path, cache_dir / slugify(pdf_path.stem), dpi=dpi, page_indices=page_indices, force=force_render)  # type: ignore[arg-type]
    for j, (idx, img) in enumerate(zip(page_indices, rendered), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_title(slide, f"{case.name}: {title_prefix}", f"{section_label} page {idx + 1}")
        add_image_contain(slide, img, 0.55, 0.92, 12.25, 5.95, border=True)
        add_footer(slide, case.name, rel_path(pdf_path, root))
    return len(rendered)


def prepare_spotlight_file_items(
    files: Sequence[Path],
    case: CaseAssets,
    cache_dir: Path,
    dpi: int,
    force_render: bool,
) -> List[Tuple[Path, Path, str]]:
    items: List[Tuple[Path, Path, str]] = []
    for f in files:
        try:
            if f.suffix.lower() == ".pdf":
                rendered = render_pdf_pages(f, cache_dir / "spotlight_files" / slugify(f.stem), dpi=dpi, max_pages=1, force=force_render)
                if rendered:
                    items.append((rendered[0], f, short_label_from_path(f, case.path)))
            else:
                img = ensure_pptx_image(f, cache_dir / "spotlight_files" / "converted", force=force_render)
                items.append((img, f, short_label_from_path(f, case.path)))
        except Exception as exc:
            print(f"WARNING: could not use spotlight file {f}: {exc}", file=sys.stderr)
    return items


def add_spotlight_grid_slides(
    prs: Presentation,
    case: CaseAssets,
    root: Path,
    items: Sequence[Tuple[Path, Path, str]],
    per_slide: int,
) -> Tuple[int, int]:
    if not items:
        return 0, 0
    per_slide = max(1, min(4, int(per_slide)))

    layouts = {
        1: [(0.65, 1.02, 12.05, 5.78)],
        2: [(0.55, 1.10, 6.05, 5.55), (6.73, 1.10, 6.05, 5.55)],
        3: [(0.55, 1.03, 6.05, 2.72), (6.73, 1.03, 6.05, 2.72), (0.55, 4.05, 6.05, 2.72)],
        4: [(0.55, 1.03, 6.05, 2.72), (6.73, 1.03, 6.05, 2.72), (0.55, 4.05, 6.05, 2.72), (6.73, 4.05, 6.05, 2.72)],
    }

    slides = 0
    for chunk_idx, chunk in enumerate(chunked(list(items), per_slide), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_title(slide, f"{case.name}: class-specific whole-slide H&E spotlights", f"spotlights {chunk_idx}")
        boxes = layouts[len(chunk)]
        source_paths: List[str] = []
        for (img, original, label), (x, y, w, h) in zip(chunk, boxes):
            add_image_contain(slide, img, x, y + 0.18, w, h - 0.18, border=True)
            add_text(slide, label, x, y - 0.02, w, 0.18, font_size=7.4, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
            source_paths.append(rel_path(original, root))
        source_txt = "; ".join(source_paths[:3])
        if len(source_paths) > 3:
            source_txt += f"; +{len(source_paths) - 3} more"
        add_footer(slide, case.name, source_txt)
        slides += 1
    return slides, len(items)


def max_pages_arg(value: int) -> Optional[int]:
    return None if value <= 0 else value


def included_regular_page_indices(pdf_path: Optional[Path], max_pages_value: int) -> set[int]:
    n = pdf_page_count(pdf_path)
    if n <= 0:
        return set()
    if max_pages_value <= 0:
        return set(range(n))
    return set(range(min(n, max_pages_value)))


def build_deck(cases: Sequence[CaseAssets], root: Path, out_path: Path, args: argparse.Namespace) -> BuildStats:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    stats = BuildStats(out=out_path, cases=len(cases))

    if not args.skip_title_slide:
        add_title_slide(prs, cases, root, include_umap=args.include_umap, compact=args.compact_first_pages, include_spotlights=args.include_spotlights)
    if not args.skip_dataset_summary:
        add_dataset_summary_slide(prs, cases, root, args)
    if not args.skip_aggregate_counts:
        stats.count_slides += add_aggregate_counts_slide(prs, cases, args.cache_dir, top_classes=args.top_classes)

    for idx, case in enumerate(cases, start=1):
        sample_cache = args.cache_dir / slugify(case.name)
        record: Dict[str, Any] = {
            "sample": case.name,
            "counts_slide": 0,
            "overview_slide": 0,
            "roi_slides": 0,
            "umap_slides": 0,
            "spotlight_slides": 0,
            "spotlight_items": 0,
            "spotlight_files": [],
            "roi_spotlight_page_indices_1based": [],
        }

        if not args.skip_section_slides:
            add_section_slide(prs, case, idx, len(cases))

        if not args.skip_counts:
            n = add_case_counts_slide(prs, case, root, sample_cache, top_classes=args.top_classes)
            stats.count_slides += n
            record["counts_slide"] = n
            if n == 0:
                stats.missing_counts.append(case.name)

        if not args.skip_overview:
            n = add_overview_slide(prs, case, root, sample_cache, dpi=args.dpi, overview_source=args.overview_source, force_render=args.force_render)
            stats.overview_slides += n
            record["overview_slide"] = n
            if n == 0:
                stats.missing_overview.append(case.name)

        n_roi = add_pdf_report_slides(
            prs,
            case,
            case.roi_pdf,
            root,
            sample_cache,
            dpi=args.dpi,
            max_pages=max_pages_arg(args.max_roi_pages),
            force_render=args.force_render,
            title_prefix="spatial ROI report",
            section_label="ROI",
        )
        stats.roi_slides += n_roi
        record["roi_slides"] = n_roi
        if n_roi == 0:
            stats.missing_roi.append(case.name)

        if args.include_spotlights:
            spotlight_files = discover_spotlight_files(
                case,
                terms=args.spotlight_terms,
                extra_globs=args.spotlight_globs,
                include_all_spatial_roi_images=args.include_all_spatial_roi_images,
                max_files=args.max_spotlight_pages,
            )
            spotlight_items = prepare_spotlight_file_items(
                spotlight_files,
                case,
                sample_cache,
                dpi=args.dpi,
                force_render=args.force_render,
            )
            file_spotlight_slides, file_spotlight_items = add_spotlight_grid_slides(
                prs,
                case,
                root,
                spotlight_items,
                per_slide=args.spotlights_per_slide,
            )
            stats.spotlight_slides += file_spotlight_slides
            stats.spotlight_items += file_spotlight_items
            stats.spotlight_files += len(spotlight_files)
            record["spotlight_slides"] += file_spotlight_slides
            record["spotlight_items"] += file_spotlight_items
            record["spotlight_files"] = [rel_path(f, root) for f in spotlight_files]

            excluded_roi_pages = included_regular_page_indices(case.roi_pdf, args.max_roi_pages)
            roi_spotlight_indices = find_pdf_pages_by_terms(
                case.roi_pdf,
                terms=args.spotlight_terms,
                exclude_indices=excluded_roi_pages,
                max_pages=args.max_roi_spotlight_pages,
            )
            if roi_spotlight_indices:
                n_roi_spot = add_pdf_selected_page_slides(
                    prs,
                    case,
                    case.roi_pdf,
                    page_indices=roi_spotlight_indices,
                    root=root,
                    cache_dir=sample_cache / "roi_spotlight_pages",
                    dpi=args.dpi,
                    force_render=args.force_render,
                    title_prefix="spatial ROI report spotlight page",
                    section_label="ROI spotlight",
                )
                stats.spotlight_slides += n_roi_spot
                stats.spotlight_items += n_roi_spot
                stats.roi_spotlight_pages += n_roi_spot
                record["spotlight_slides"] += n_roi_spot
                record["spotlight_items"] += n_roi_spot
                record["roi_spotlight_page_indices_1based"] = [i + 1 for i in roi_spotlight_indices]

            if not spotlight_files and not roi_spotlight_indices:
                stats.missing_spotlights.append(case.name)

        if args.include_umap:
            n_umap = add_pdf_report_slides(
                prs,
                case,
                case.umap_pdf,
                root,
                sample_cache,
                dpi=args.dpi,
                max_pages=max_pages_arg(args.max_umap_pages),
                force_render=args.force_render,
                title_prefix="cell-embedding UMAP report",
                section_label="UMAP",
            )
            stats.umap_slides += n_umap
            record["umap_slides"] = n_umap

        stats.case_records.append(record)

    stats.slides = len(prs.slides)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return stats


def estimate_case_slides(case: CaseAssets, args: argparse.Namespace) -> Dict[str, Any]:
    out: Dict[str, Any] = {}
    out["section"] = 0 if args.skip_section_slides else 1
    out["counts"] = 0 if args.skip_counts else (1 if safe_exists(case.counts_csv) else 0)
    out["overview"] = 0 if args.skip_overview else (1 if (safe_exists(case.overview_pdf) or safe_exists(case.overview_png)) else 0)
    roi_pages = pdf_page_count(case.roi_pdf)
    if roi_pages > 0:
        out["roi"] = roi_pages if args.max_roi_pages <= 0 else min(roi_pages, args.max_roi_pages)
    else:
        out["roi"] = 0
    umap_pages = pdf_page_count(case.umap_pdf)
    if args.include_umap and umap_pages > 0:
        out["umap"] = umap_pages if args.max_umap_pages <= 0 else min(umap_pages, args.max_umap_pages)
    else:
        out["umap"] = 0
    if args.include_spotlights:
        spotlight_files = discover_spotlight_files(
            case,
            terms=args.spotlight_terms,
            extra_globs=args.spotlight_globs,
            include_all_spatial_roi_images=args.include_all_spatial_roi_images,
            max_files=args.max_spotlight_pages,
        )
        out["spotlight_files"] = len(spotlight_files)
        out["spotlight_file_slides_est"] = int(math.ceil(len(spotlight_files) / max(1, args.spotlights_per_slide)))
        excluded = included_regular_page_indices(case.roi_pdf, args.max_roi_pages)
        roi_spot_indices = find_pdf_pages_by_terms(case.roi_pdf, args.spotlight_terms, excluded, args.max_roi_spotlight_pages)
        out["roi_spotlight_pages"] = len(roi_spot_indices)
        out["spotlight"] = out["spotlight_file_slides_est"] + out["roi_spotlight_pages"]
    else:
        out["spotlight_files"] = "off"
        out["roi_spotlight_pages"] = "off"
        out["spotlight"] = 0
    out["total_est"] = out["section"] + out["counts"] + out["overview"] + out["roi"] + out["umap"] + out["spotlight"]
    return out


def print_dry_run(cases: Sequence[CaseAssets], root: Path, args: argparse.Namespace) -> None:
    print(f"Script: {Path(__file__).name} | {SCRIPT_VERSION}")
    print(f"Detected samples: {len(cases)}")
    print(f"Output PPTX: {args.out}")
    if not args.no_cell_stats_csv:
        print(f"Cell stats raw CSV: {args.cell_stats_csv}")
        print(f"Cell stats normalized CSV: {args.cell_stats_normalized_csv}")
    else:
        print("Cell stats CSV outputs: disabled")
    print(f"DPI: {args.dpi}")
    print(f"Compact first pages: {args.compact_first_pages}")
    print(f"Include UMAP: {args.include_umap}")
    print(f"Include spotlights: {args.include_spotlights}")
    print(f"Max ROI pages per sample: {'all' if args.max_roi_pages <= 0 else args.max_roi_pages}")
    print(f"Max UMAP pages per sample: {'all' if args.max_umap_pages <= 0 else args.max_umap_pages}")
    if args.include_spotlights:
        print(f"Max standalone spotlight files per sample: {'all' if args.max_spotlight_pages <= 0 else args.max_spotlight_pages}")
        print(f"Max ROI-report spotlight pages per sample: {'all' if args.max_roi_spotlight_pages <= 0 else args.max_roi_spotlight_pages}")
        print(f"Spotlights per slide: {args.spotlights_per_slide}")
    print("")

    total_est = 0
    for case in cases:
        cells = read_case_cell_total(case)
        est = estimate_case_slides(case, args)
        total_est += int(est["total_est"])
        print(f"[{case.name}]")
        print(f"  cells:                  {cells if cells is not None else 'NA'}")
        print(f"  counts_csv:             {rel_path(case.counts_csv, root)}")
        print(f"  overview_pdf:           {rel_path(case.overview_pdf, root)}")
        print(f"  overview_png:           {rel_path(case.overview_png, root)}")
        print(f"  roi_pdf:                {rel_path(case.roi_pdf, root)}  pages={pdf_page_count(case.roi_pdf)}  included={est['roi']}")
        print(f"  umap_pdf:               {rel_path(case.umap_pdf, root)}  pages={pdf_page_count(case.umap_pdf)}  included={est['umap']}")
        print(f"  spotlight_files:        {est['spotlight_files']}")
        print(f"  roi_spotlight_pages:    {est['roi_spotlight_pages']}")
        print(f"  estimated sample slides:{est['total_est']}")
        print("")

    overhead = (0 if args.skip_title_slide else 1) + (0 if args.skip_dataset_summary else 1)
    if not args.skip_aggregate_counts and aggregate_counts(cases) is not None:
        overhead += 1
    print(f"Estimated total slides: {total_est + overhead}  (sample slides {total_est} + overhead {overhead})")


def parse_args(argv: Optional[Sequence[str]] = None) -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="Create ONE cancer-agnostic combined PPTX from LazySlide/HistoPLUS sample folders, plus companion sample x cell-type CSV matrices.",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    p.add_argument("--version", action="version", version=f"%(prog)s {SCRIPT_VERSION}")
    p.add_argument("--root", required=True, type=Path, help="Root folder containing one folder per cancer sample/case.")
    p.add_argument("--out", required=True, type=Path, help="Output .pptx path for the single combined deck.")
    p.add_argument("--recursive", action="store_true", help="Search sample folders recursively under --root.")
    p.add_argument("--include", nargs="*", default=[], help="Optional shell-style sample name filters, e.g. '*HE_1' '4904*'.")
    p.add_argument("--exclude", nargs="*", default=[], help="Optional shell-style sample name exclusions.")
    p.add_argument("--first-n", type=int, default=0, help="Debug/test: only include the first N detected samples. Use 0 for all.")

    p.add_argument("--dpi", type=int, default=180, help="DPI used to rasterize PDF pages before placing them in PPTX.")
    p.add_argument("--overview-source", choices=["auto", "pdf", "png"], default="auto", help="Source for overview slide. auto prefers PDF then PNG.")
    p.add_argument("--max-roi-pages", type=int, default=0, help="Max pages from each spatial_rois_report.pdf. Use 0 for all pages. In --compact-first-pages, 0 becomes 1.")
    p.add_argument("--include-umap", action="store_true", help="Also include post_visualization/cell_embedding_umap_report.pdf pages.")
    p.add_argument("--max-umap-pages", type=int, default=0, help="Max pages from each UMAP PDF. Use 0 for all pages. In --compact-first-pages, 0 becomes 1.")
    p.add_argument("--top-classes", type=int, default=15, help="Number of classes shown in count barplots.")

    p.add_argument("--cell-stats-csv", type=Path, default=None, help="Output CSV for raw detected cell counts. Rows=samples, columns=cell types/classes. Default: <out>.cell_counts_by_sample.csv.")
    p.add_argument("--cell-stats-normalized-csv", type=Path, default=None, help="Output CSV for row-normalized cell fractions. Rows=samples, columns=cell types/classes. Default: <out>.cell_fractions_by_sample.csv.")
    p.add_argument("--no-cell-stats-csv", action="store_true", help="Disable the two companion sample x cell-type CSV files.")

    p.add_argument("--compact-first-pages", action="store_true", help="Compact mode: first page of ROI report and first UMAP page only; also enables spotlight search unless --no-spotlights is passed.")
    p.add_argument("--include-spotlights", action="store_true", help="Include class-specific whole-slide H&E spotlight files/pages when detected.")
    p.add_argument("--no-spotlights", action="store_true", help="Disable spotlights, useful with --compact-first-pages when you want first report pages only.")
    p.add_argument("--max-spotlight-pages", type=int, default=12, help="Max standalone spotlight image/PDF files per sample. Use 0 for all.")
    p.add_argument("--max-roi-spotlight-pages", type=int, default=6, help="Max additional spotlight-like pages found inside spatial_rois_report.pdf per sample. Use 0 for all.")
    p.add_argument("--spotlights-per-slide", type=int, default=4, help="Number of standalone spotlight images placed per slide. Use 1, 2, 3, or 4.")
    p.add_argument("--spotlight-terms", nargs="*", default=list(DEFAULT_SPOTLIGHT_TERMS), help="Terms used to detect spotlight files and spotlight pages in PDFs.")
    p.add_argument("--spotlight-globs", nargs="*", default=[], help="Extra shell globs relative to each sample folder, e.g. 'post_visualization/spatial_rois/**/*spotlight*.png'.")
    p.add_argument("--include-all-spatial-roi-images", action="store_true", help="Broad mode: treat all image/PDF files under post_visualization/spatial_rois as spotlight candidates; combine with --max-spotlight-pages.")

    p.add_argument("--skip-title-slide", action="store_true", help="Skip the first title slide.")
    p.add_argument("--skip-dataset-summary", action="store_true", help="Skip the dataset/sample availability table slide.")
    p.add_argument("--skip-section-slides", action="store_true", help="Skip per-sample divider slides to make the deck smaller.")
    p.add_argument("--skip-counts", action="store_true", help="Skip per-sample class-count slides.")
    p.add_argument("--skip-aggregate-counts", action="store_true", help="Skip aggregate all-sample count slide.")
    p.add_argument("--skip-overview", action="store_true", help="Skip celltype overview/zoom slides.")

    p.add_argument("--dry-run", action="store_true", help="Print detected samples/files and estimated slides; exit without creating PPTX.")
    p.add_argument("--force-render", action="store_true", help="Force regeneration of cached PNGs from PDFs/images.")
    p.add_argument("--clean-cache", action="store_true", help="Delete render cache before starting.")
    p.add_argument("--cache-dir", type=Path, default=None, help="Optional cache directory. Default: <outdir>/_render_cache_all_samples_V3.")

    args = p.parse_args(argv)
    args.root = args.root.expanduser().resolve()
    args.out = args.out.expanduser().resolve()
    if args.cache_dir is None:
        args.cache_dir = args.out.parent / "_render_cache_all_samples_V3"
    else:
        args.cache_dir = args.cache_dir.expanduser().resolve()

    if args.cell_stats_csv is None:
        args.cell_stats_csv = args.out.with_suffix(".cell_counts_by_sample.csv")
    else:
        args.cell_stats_csv = args.cell_stats_csv.expanduser().resolve()
    if args.cell_stats_normalized_csv is None:
        args.cell_stats_normalized_csv = args.out.with_suffix(".cell_fractions_by_sample.csv")
    else:
        args.cell_stats_normalized_csv = args.cell_stats_normalized_csv.expanduser().resolve()

    if args.compact_first_pages:
        if args.max_roi_pages <= 0:
            args.max_roi_pages = 1
        if args.max_umap_pages <= 0:
            args.max_umap_pages = 1
        if not args.no_spotlights:
            args.include_spotlights = True

    if args.no_spotlights and args.include_spotlights:
        p.error("Use either --include-spotlights or --no-spotlights, not both.")
    if not args.root.exists() or not args.root.is_dir():
        p.error(f"--root does not exist or is not a directory: {args.root}")
    if args.out.suffix.lower() != ".pptx":
        p.error("--out must end in .pptx")
    if args.dpi < 72 or args.dpi > 400:
        p.error("--dpi must be between 72 and 400. Recommended: 150-220.")
    if args.top_classes < 1:
        p.error("--top-classes must be >= 1")
    if args.first_n < 0:
        p.error("--first-n must be >= 0")
    if args.max_roi_pages < 0:
        p.error("--max-roi-pages must be >= 0")
    if args.max_umap_pages < 0:
        p.error("--max-umap-pages must be >= 0")
    if args.max_spotlight_pages < 0:
        p.error("--max-spotlight-pages must be >= 0")
    if args.max_roi_spotlight_pages < 0:
        p.error("--max-roi-spotlight-pages must be >= 0")
    if args.spotlights_per_slide < 1 or args.spotlights_per_slide > 4:
        p.error("--spotlights-per-slide must be 1, 2, 3, or 4")

    return args


def main(argv: Optional[Sequence[str]] = None) -> int:
    args = parse_args(argv)
    args.out.parent.mkdir(parents=True, exist_ok=True)

    if args.clean_cache and args.cache_dir.exists():
        # Safety guard: never delete root or output directory by accident.
        if args.cache_dir in {args.root, args.out.parent, Path("/")}:
            print(f"ERROR: refusing to delete unsafe cache dir: {args.cache_dir}", file=sys.stderr)
            return 2
        shutil.rmtree(args.cache_dir)
    args.cache_dir.mkdir(parents=True, exist_ok=True)

    cases = discover_cases(args.root, recursive=args.recursive, include=args.include, exclude=args.exclude, first_n=args.first_n)
    if not cases:
        print("ERROR: no LazySlide/HistoPLUS sample folders were detected.", file=sys.stderr)
        print("Check --root and expected files under overlays/, cell_types/, and post_visualization/.", file=sys.stderr)
        return 1

    if args.dry_run:
        print_dry_run(cases, args.root, args)
        return 0

    cell_stats_info: Dict[str, Any] = {}
    if not args.no_cell_stats_csv:
        cell_stats_info = write_cell_stats_csvs(
            cases,
            raw_csv=args.cell_stats_csv,
            normalized_csv=args.cell_stats_normalized_csv,
        )
        print(f"Wrote cell-count matrix CSV: {args.cell_stats_csv}")
        print(f"Wrote normalized cell-fraction matrix CSV: {args.cell_stats_normalized_csv}")
        print(f"  detected cell types/classes: {cell_stats_info.get('detected_cell_types', 0)}")
        print(f"  total detected cells: {cell_stats_info.get('total_detected_cells', 0):,}")

    stats = build_deck(cases, args.root, args.out, args)
    print(f"Wrote single combined PPTX: {args.out}")
    print(f"  script: {Path(__file__).name}")
    print(f"  version: {SCRIPT_VERSION}")
    print(f"  samples: {stats.cases}")
    print(f"  slides: {stats.slides}")
    print(f"  overview slides: {stats.overview_slides}")
    print(f"  count slides: {stats.count_slides}")
    print(f"  ROI slides: {stats.roi_slides}")
    print(f"  UMAP slides: {stats.umap_slides}")
    print(f"  spotlight slides: {stats.spotlight_slides}")
    print(f"  spotlight items: {stats.spotlight_items}")
    print(f"  standalone spotlight files detected: {stats.spotlight_files}")
    print(f"  ROI-report spotlight pages added: {stats.roi_spotlight_pages}")

    manifest = {
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "script": Path(__file__).name,
        "script_version": SCRIPT_VERSION,
        "root": str(args.root),
        "out": str(args.out),
        "cache_dir": str(args.cache_dir),
        "n_cases": len(cases),
        "cases": [c.name for c in cases],
        "options": {
            "dpi": args.dpi,
            "overview_source": args.overview_source,
            "compact_first_pages": args.compact_first_pages,
            "max_roi_pages": args.max_roi_pages,
            "include_umap": args.include_umap,
            "max_umap_pages": args.max_umap_pages,
            "top_classes": args.top_classes,
            "cell_stats_csv_enabled": not args.no_cell_stats_csv,
            "cell_stats_csv": str(args.cell_stats_csv) if not args.no_cell_stats_csv else None,
            "cell_stats_normalized_csv": str(args.cell_stats_normalized_csv) if not args.no_cell_stats_csv else None,
            "include_spotlights": args.include_spotlights,
            "max_spotlight_pages": args.max_spotlight_pages,
            "max_roi_spotlight_pages": args.max_roi_spotlight_pages,
            "spotlights_per_slide": args.spotlights_per_slide,
            "spotlight_terms": args.spotlight_terms,
            "spotlight_globs": args.spotlight_globs,
            "include_all_spatial_roi_images": args.include_all_spatial_roi_images,
            "skip_title_slide": args.skip_title_slide,
            "skip_dataset_summary": args.skip_dataset_summary,
            "skip_section_slides": args.skip_section_slides,
            "skip_counts": args.skip_counts,
            "skip_aggregate_counts": args.skip_aggregate_counts,
            "skip_overview": args.skip_overview,
        },
        "cell_stats": cell_stats_info,
        "stats": stats.as_dict(),
    }
    manifest_path = args.out.with_suffix(".manifest.json")
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)
    print(f"Wrote manifest: {manifest_path}")
    if args.include_spotlights and stats.spotlight_items == 0:
        print("WARNING: --include-spotlights was active, but no spotlight files/pages were detected.", file=sys.stderr)
        print("         Try --include-all-spatial-roi-images or add --spotlight-globs 'post_visualization/spatial_rois/**/*.png'", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
